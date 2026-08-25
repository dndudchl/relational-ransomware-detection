#!/usr/bin/env python3
"""
make_documents.py - Generate document files to submit as benign samples.

Why documents
-------------
The hard negatives built so far have a weakness a reviewer will find: most of
them were written for this experiment. Wrappers around 7-Zip, cipher and
robocopy answer that partly, but they are 24 of 68.

Submitting a .docx does something different. The sandbox opens it with
whatever is registered to handle it, so the program under analysis is Word or
Excel or Acrobat -- signed by Microsoft or Adobe, installed on the guest
before this project began, and doing the thing it exists to do. The file
activity that produces is real: the document is read, temporary files appear
beside it, the lock file ~$name.docx is created and removed, the recent items
list is rewritten, autorecovery data is saved.

That last part matters, because it is the shape that caused a false positive
earlier. tool_open_docs opened twelve decoy documents and was classified as
ransomware with a score of 1.00. Submitting documents directly turns that one
observation into a population.

What varies
-----------
Size, structure and type, because a one-paragraph document and a
two-hundred-page one make different amounts of work for the handler, and the
question is where along that range the detector starts firing. Content is
irrelevant -- the handler does the same thing whatever the words are -- so it
is filler.

Usage
-----
  python3 make_documents.py --count 300 --outdir ~/documents
"""

import os
import csv
import random
import argparse

FILLER = """The quarterly review covers operational performance across the
period, with attention to throughput, error rates and the outstanding items
carried forward from the previous cycle. Figures are provisional until the
reconciliation completes. Departmental submissions were received on time with
two exceptions, both noted in the appendix. No material changes to the
forecast are proposed at this stage.""".replace("\n", " ")


def make_docx(path, paragraphs, tables):
    from docx import Document
    from docx.shared import Pt
    doc = Document()
    doc.add_heading("Operational Review", level=1)
    for i in range(paragraphs):
        if i % 8 == 0:
            doc.add_heading(f"Section {i // 8 + 1}", level=2)
        p = doc.add_paragraph(FILLER)
        p.runs[0].font.size = Pt(11)
    for t in range(tables):
        table = doc.add_table(rows=6, cols=4)
        table.style = "Table Grid"
        for r in range(6):
            for c in range(4):
                table.cell(r, c).text = f"{r * 4 + c}"
    doc.save(path)


def make_xlsx(path, sheets, rows):
    from openpyxl import Workbook
    wb = Workbook()
    for s in range(sheets):
        ws = wb.active if s == 0 else wb.create_sheet()
        ws.title = f"Data{s + 1}"
        ws.append(["id", "region", "units", "value", "note"])
        for r in range(rows):
            ws.append([r, f"R{r % 7}", r * 3, round(r * 1.7, 2), "ok"])
        # A formula gives the handler something to evaluate on open, which is
        # work a plain table does not produce.
        ws.cell(row=rows + 2, column=4, value=f"=SUM(D2:D{rows + 1})")
    wb.save(path)


def make_pptx(path, slides):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    for i in range(slides):
        layout = prs.slide_layouts[1 if i else 0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"Section {i + 1}"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = FILLER[:300]
    prs.save(path)


def make_csv(path, rows):
    with open(path, "w", newline="") as f:
        w = csv.writer(f)
        w.writerow(["id", "region", "units", "value", "note"])
        for r in range(rows):
            w.writerow([r, f"R{r % 7}", r * 3, round(r * 1.7, 2), "ok"])


def make_rtf(path, paragraphs):
    body = "".join("\\par " + FILLER for _ in range(paragraphs))
    with open(path, "w") as f:
        f.write(r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Times New Roman;}}" +
                r"\f0\fs22 " + body + "}")


def make_txt(path, paragraphs):
    with open(path, "w") as f:
        for _ in range(paragraphs):
            f.write(FILLER + "\n\n")


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    # Eighty rather than several hundred. Word opening the fiftieth
    # two-hundred-paragraph document does the same thing it did on the tenth,
    # so the information is in the grid of format against size -- six formats
    # by three or four size steps -- and a few instances of each cell is
    # enough to see the variance within it. The remaining analysis budget is
    # better spent on variants and installers, which are not interchangeable
    # with each other.
    parser.add_argument("--count", type=int, default=80)
    parser.add_argument("--outdir", default="./documents")
    parser.add_argument("--seed", type=int, default=0)
    args = parser.parse_args()

    outdir = os.path.expanduser(args.outdir)
    os.makedirs(outdir, exist_ok=True)
    rng = random.Random(args.seed)

    missing = []
    for mod, pkg in [("docx", "python-docx"), ("openpyxl", "openpyxl"),
                     ("pptx", "python-pptx")]:
        try:
            __import__(mod)
        except ImportError:
            missing.append(pkg)
    if missing:
        print("[!] install first:  pip install --break-system-packages "
              + " ".join(missing))
        return

    # The grid, enumerated rather than sampled.
    #
    # Drawing format and size at random leaves the coverage to luck: a run of
    # 25 came out with eleven spreadsheets and five documents, which is not
    # enough Word files to cover the four sizes. There are only twenty
    # combinations, so cycling through them in order covers every one before
    # any repeats, and the count decides how many times round.
    GRID = (
        [("docx", (p, t)) for p, t in [(3, 0), (20, 1), (80, 3), (200, 6)]]
        + [("xlsx", (sh, r)) for sh, r in [(1, 20), (1, 500), (3, 200), (5, 2000)]]
        + [("pptx", (n,)) for n in (3, 10, 30)]
        + [("csv", (n,)) for n in (50, 1000, 20000)]
        + [("rtf", (n,)) for n in (5, 40, 150)]
        + [("txt", (n,)) for n in (10, 100, 800)]
    )

    rows = []
    for i in range(args.count):
        kind, spec = GRID[i % len(GRID)]
        name = f"doc_{i:04d}.{kind}"
        path = os.path.join(outdir, name)
        try:
            if kind == "docx":
                paras, tables = spec
                make_docx(path, paras, tables)
                size = f"{paras} paragraphs, {tables} tables"
            elif kind == "xlsx":
                sheets, r = spec
                make_xlsx(path, sheets, r)
                size = f"{sheets} sheets, {r} rows"
            elif kind == "pptx":
                make_pptx(path, spec[0])
                size = f"{spec[0]} slides"
            elif kind == "csv":
                make_csv(path, spec[0])
                size = f"{spec[0]} rows"
            elif kind == "rtf":
                make_rtf(path, spec[0])
                size = f"{spec[0]} paragraphs"
            else:
                make_txt(path, spec[0])
                size = f"{spec[0]} paragraphs"
        except Exception as e:
            print(f"[!] {name}: {type(e).__name__}: {e}")
            continue
        rows.append({"filename": name, "kind": kind, "shape": size,
                     "bytes": os.path.getsize(path)})
        if (i + 1) % 50 == 0:
            print(f"\r   {i + 1}/{args.count}", end="", flush=True)
    print()

    manifest = os.path.join(outdir, "document_manifest.csv")
    with open(manifest, "w", newline="") as f:
        w = csv.DictWriter(f, fieldnames=["filename", "kind", "shape", "bytes"])
        w.writeheader(); w.writerows(rows)

    from collections import Counter
    total = sum(r["bytes"] for r in rows)
    print(f"{len(rows)} documents, {total / 1e6:.1f} MB")
    for k, n in sorted(Counter(r["kind"] for r in rows).items()):
        print(f"   {k:<6}{n:>5}")
    shapes = Counter((r["kind"], r["shape"]) for r in rows)
    print(f"   {len(shapes)} of {len(GRID)} combinations covered")
    print(f"[saved] {manifest}")
    print("\nSubmit these the same way as everything else. The program the")
    print("sandbox records will be the registered handler, not the file.")


if __name__ == "__main__":
    main()
